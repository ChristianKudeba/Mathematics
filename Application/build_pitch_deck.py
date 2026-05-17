"""Generate the mathAI vision deck (power-user aesthetic, dark theme)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- palette ----------
BG       = RGBColor(0x0B, 0x0E, 0x12)   # near black
PANEL    = RGBColor(0x12, 0x17, 0x1F)   # subtle panel
INK      = RGBColor(0xE6, 0xEA, 0xEF)   # body
DIM      = RGBColor(0x8A, 0x95, 0xA5)   # secondary
RULE     = RGBColor(0x2A, 0x33, 0x40)   # hairline
ACCENT   = RGBColor(0x00, 0xE5, 0xFF)   # electric cyan
AMBER    = RGBColor(0xFF, 0xB3, 0x47)   # warm amber
GREEN    = RGBColor(0x7C, 0xE3, 0xA8)   # verified
MAGENTA  = RGBColor(0xFF, 0x6E, 0xC7)

FONT_BODY = "Segoe UI"
FONT_MONO = "Consolas"

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ---------- helpers ----------
def add_bg(slide, color=BG):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid(); rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    return rect

def add_rect(slide, x, y, w, h, color, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s

def add_text(slide, x, y, w, h, text, *, font=FONT_BODY, size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, tracking=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        if tracking:
            # letter-spacing via raw XML (spc is in units of 1/100 pt)
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(int(tracking * 100)))
    return tb

def add_chrome(slide, eyebrow, slide_num, total):
    # top hairline
    add_rect(slide, Inches(0.6), Inches(0.55), Inches(12.13), Emu(9525), RULE)
    # eyebrow upper-left
    add_text(slide, Inches(0.6), Inches(0.22), Inches(8), Inches(0.3),
             eyebrow.upper(), font=FONT_MONO, size=10, color=ACCENT, tracking=2)
    # slide counter upper-right
    add_text(slide, Inches(11.0), Inches(0.22), Inches(1.73), Inches(0.3),
             f"{slide_num:02d} / {total:02d}", font=FONT_MONO, size=10,
             color=DIM, align=PP_ALIGN.RIGHT, tracking=1)
    # bottom hairline
    add_rect(slide, Inches(0.6), Inches(7.0), Inches(12.13), Emu(9525), RULE)
    # footer
    add_text(slide, Inches(0.6), Inches(7.10), Inches(8), Inches(0.3),
             "MATHAI · CONFIDENTIAL DRAFT · 2026.05",
             font=FONT_MONO, size=9, color=DIM, tracking=2)
    add_text(slide, Inches(11.0), Inches(7.10), Inches(1.73), Inches(0.3),
             "// vision deck", font=FONT_MONO, size=9, color=DIM,
             align=PP_ALIGN.RIGHT, tracking=1)


def section_label(slide, text, x=Inches(0.6), y=Inches(0.9)):
    # vertical accent bar + label
    add_rect(slide, x, y, Emu(38100), Inches(0.32), ACCENT)  # ~3pt wide
    add_text(slide, x + Inches(0.18), y - Inches(0.02), Inches(8), Inches(0.36),
             text.upper(), font=FONT_MONO, size=11, color=ACCENT, tracking=3)


TOTAL = 13

# ============================================================
# SLIDE 1 — Title
# ============================================================
s = prs.slides.add_slide(blank); add_bg(s)
add_chrome(s, "// 00 cover", 1, TOTAL)

# ambient grid lines (subtle)
for i in range(1, 12):
    add_rect(s, Inches(0.6 + i * 1.0), Inches(0.6), Emu(6350), Inches(6.4), RULE)

# vertical accent bar
add_rect(s, Inches(0.6), Inches(2.7), Emu(76200), Inches(2.4), ACCENT)

add_text(s, Inches(0.95), Inches(2.55), Inches(8), Inches(0.4),
         "MATHAI", font=FONT_MONO, size=12, color=ACCENT, tracking=6)

add_text(s, Inches(0.95), Inches(2.95), Inches(11.5), Inches(1.6),
         "A living map of\nmathematics.",
         font=FONT_BODY, size=72, bold=True, color=INK)

add_text(s, Inches(0.95), Inches(5.05), Inches(11.5), Inches(0.6),
         "Autonomous research agents for the people who actually do the math.",
         font=FONT_BODY, size=22, color=DIM)

add_text(s, Inches(0.95), Inches(6.15), Inches(11.5), Inches(0.4),
         "Anton Shakov  ·  Christian Kudeba",
         font=FONT_MONO, size=11, color=AMBER, tracking=2)


# ============================================================
# SLIDE 2 — Vision (the manifesto)
# ============================================================
s = prs.slides.add_slide(blank); add_bg(s)
add_chrome(s, "// 01 vision", 2, TOTAL)
section_label(s, "Vision")

add_text(s, Inches(0.6), Inches(1.35), Inches(11.5), Inches(0.8),
         "What we are building.",
         font=FONT_BODY, size=40, bold=True, color=INK)

# three columns of paragraphs
col_w = Inches(3.85)
gap   = Inches(0.25)
x0    = Inches(0.6)
y0    = Inches(2.55)
col_h = Inches(4.2)

paragraphs = [
    ("01", "THE AGENT",
     "For four centuries, mathematics has been done by hand — one mind, one notebook, "
     "one proof at a time. We change the unit of work. Each researcher shapes a "
     "persistent AI agent that thinks the way they think, points at the problems "
     "they care about, and runs for hours, days, weeks. Their agent. Their name "
     "on what it finds."),
    ("02", "THE ARTIFACT",
     "Every result that survives verification — every theorem, counterexample, "
     "conjecture-with-evidence — flows into a single growing public artifact: "
     "a mechanically-checked, attribution-preserving map of mathematics. Free to "
     "read. Owned by no one. Stewarded by the field itself."),
    ("03", "THE MOMENT",
     "The conditions only just arrived. Frontier models reason at graduate level "
     "when steered well. Verification has become tractable at scale. The cost of "
     "letting an agent think for a week has collapsed. And the field's most "
     "respected names have, for the first time, said the quiet part out loud: "
     "this is ready."),
]

for i, (num, head, body) in enumerate(paragraphs):
    x = x0 + i * (col_w + gap)
    # accent rule on top of column
    add_rect(s, x, y0, col_w, Emu(19050), ACCENT)
    add_text(s, x, y0 + Inches(0.18), col_w, Inches(0.32),
             num, font=FONT_MONO, size=11, color=ACCENT, tracking=2)
    add_text(s, x, y0 + Inches(0.55), col_w, Inches(0.4),
             head, font=FONT_MONO, size=12, bold=True, color=AMBER, tracking=2)
    add_text(s, x, y0 + Inches(1.05), col_w, col_h,
             body, font=FONT_BODY, size=14, color=INK)


# ============================================================
# SLIDE 3 — The Moment (why now)
# ============================================================
s = prs.slides.add_slide(blank); add_bg(s)
add_chrome(s, "// 02 timing", 3, TOTAL)
section_label(s, "Why now")

add_text(s, Inches(0.6), Inches(1.35), Inches(11.5), Inches(0.8),
         "Four conditions converged.",
         font=FONT_BODY, size=40, bold=True, color=INK)

cells = [
    ("01", "REASONING", "Frontier models reach graduate-research level when steered with care."),
    ("02", "VERIFICATION", "Mechanical proof-checking finally scales. The community trusts it."),
    ("03", "ECONOMICS", "Caching and tiered routing collapsed the cost of letting an agent think."),
    ("04", "PERMISSION", "The field's most respected voices have endorsed the moment publicly."),
]
gx, gy = Inches(0.6), Inches(2.65)
cw, ch = Inches(5.95), Inches(2.0)
gap = Inches(0.23)
for i, (num, head, body) in enumerate(cells):
    r, c = divmod(i, 2)
    x = gx + c * (cw + gap)
    y = gy + r * (ch + gap)
    add_rect(s, x, y, cw, ch, PANEL, line=RULE)
    add_text(s, x + Inches(0.35), y + Inches(0.25), Inches(1.5), Inches(0.3),
             num, font=FONT_MONO, size=11, color=ACCENT, tracking=2)
    add_text(s, x + Inches(0.35), y + Inches(0.55), cw - Inches(0.7), Inches(0.45),
             head, font=FONT_BODY, size=22, bold=True, color=INK, tracking=1)
    add_text(s, x + Inches(0.35), y + Inches(1.05), cw - Inches(0.7), Inches(0.85),
             body, font=FONT_BODY, size=14, color=DIM)


# ============================================================
# SLIDE 4 — The Opportunity (trust gap)
# ============================================================
s = prs.slides.add_slide(blank); add_bg(s)
add_chrome(s, "// 03 opportunity", 4, TOTAL)
section_label(s, "The gap")

add_text(s, Inches(0.6), Inches(1.35), Inches(11.5), Inches(1.4),
         "AI can do graduate-level math.\nThe field doesn't trust it yet.",
         font=FONT_BODY, size=40, bold=True, color=INK)

# big single statement panel
add_rect(s, Inches(0.6), Inches(4.0), Inches(12.13), Inches(2.5), PANEL, line=RULE)
add_text(s, Inches(1.1), Inches(4.25), Inches(11), Inches(0.4),
         "// THESIS", font=FONT_MONO, size=11, color=AMBER, tracking=3)
add_text(s, Inches(1.1), Inches(4.65), Inches(11), Inches(1.7),
         "Mathematicians are a market that pays dearly for instruments they trust\n"
         "and ignores everything else. Trust is the entire game.",
         font=FONT_BODY, size=22, color=INK)
add_text(s, Inches(1.1), Inches(6.0), Inches(11), Inches(0.4),
         "Hallucination · No persistence · No verification · No attribution",
         font=FONT_MONO, size=12, color=ACCENT, tracking=2)


# ============================================================
# SLIDE 5 — The white space (niches)
# ============================================================
s = prs.slides.add_slide(blank); add_bg(s)
add_chrome(s, "// 04 white space", 5, TOTAL)
section_label(s, "The white space")

add_text(s, Inches(0.6), Inches(1.35), Inches(11.5), Inches(0.8),
         "Three crowded categories. One missing one.",
         font=FONT_BODY, size=36, bold=True, color=INK)

# three small "competitor" panels at top
panels = [
    ("OLYMPIAD PROVERS", "Race to win benchmarks. Narrow. Well-funded. Not a workflow."),
    ("AI CODING HARNESSES", "Powerful general tools. Built for engineers, not researchers."),
    ("NON-PROFIT SCIENCE AGENTS", "Right shape, adjacent science. Donation-funded. Free."),
]
y0 = Inches(2.6); ph = Inches(1.45); pw = Inches(3.95); gap = Inches(0.23)
for i, (head, body) in enumerate(panels):
    x = Inches(0.6) + i * (pw + gap)
    add_rect(s, x, y0, pw, ph, PANEL, line=RULE)
    add_text(s, x + Inches(0.3), y0 + Inches(0.22), pw - Inches(0.6), Inches(0.4),
             head, font=FONT_MONO, size=11, bold=True, color=DIM, tracking=2)
    add_text(s, x + Inches(0.3), y0 + Inches(0.62), pw - Inches(0.6), Inches(0.8),
             body, font=FONT_BODY, size=13, color=INK)

# big our-position panel
y1 = Inches(4.45)
add_rect(s, Inches(0.6), y1, Inches(12.13), Inches(2.4), BG, line=ACCENT)
# left accent bar
add_rect(s, Inches(0.6), y1, Emu(76200), Inches(2.4), ACCENT)
add_text(s, Inches(0.95), y1 + Inches(0.25), Inches(11), Inches(0.4),
         "// OUR POSITION", font=FONT_MONO, size=11, color=ACCENT, tracking=3)
add_text(s, Inches(0.95), y1 + Inches(0.65), Inches(11.5), Inches(1.6),
         "Verifier-pluggable autonomous research agents\n"
         "producing a permanent, public, attributed record of mathematics.",
         font=FONT_BODY, size=24, bold=True, color=INK)
add_text(s, Inches(0.95), y1 + Inches(1.95), Inches(11), Inches(0.4),
         "No funded for-profit owns this position. We are first.",
         font=FONT_MONO, size=12, color=AMBER, tracking=1)


# ============================================================
# SLIDE 6 — The product (3 columns)
# ============================================================
s = prs.slides.add_slide(blank); add_bg(s)
add_chrome(s, "// 05 product", 6, TOTAL)
section_label(s, "Product")

add_text(s, Inches(0.6), Inches(1.35), Inches(11.5), Inches(0.8),
         "Three primitives. One platform.",
         font=FONT_BODY, size=40, bold=True, color=INK)

cols = [
    ("AGENTS",     ACCENT,
     "Custom autonomous reasoners.",
     "Each researcher composes prompts, models, tools, and verifiers into "
     "a bespoke agent that thinks on long horizons about the problems they own."),
    ("VERIFICATION", GREEN,
     "Mechanical proof-checking.",
     "What can be checked, is. Lean-grade verification on every public claim. "
     "Human review for what falls outside. Trust earned the way the field grants it."),
    ("THE GRAPH",  AMBER,
     "A living public artifact.",
     "Every verified result becomes a node. Edges are dependencies, "
     "generalizations, refutations. Free to read. Owned by no one. "
     "Citable to its maker."),
]
y0 = Inches(2.65); cw = Inches(3.95); ch = Inches(4.0); gap = Inches(0.23)
for i, (head, color, lede, body) in enumerate(cols):
    x = Inches(0.6) + i * (cw + gap)
    add_rect(s, x, y0, cw, ch, PANEL, line=RULE)
    # top color bar
    add_rect(s, x, y0, cw, Inches(0.08), color)
    add_text(s, x + Inches(0.3), y0 + Inches(0.4), cw - Inches(0.6), Inches(0.4),
             f"0{i+1}", font=FONT_MONO, size=11, color=color, tracking=2)
    add_text(s, x + Inches(0.3), y0 + Inches(0.75), cw - Inches(0.6), Inches(0.55),
             head, font=FONT_BODY, size=26, bold=True, color=INK, tracking=1)
    add_text(s, x + Inches(0.3), y0 + Inches(1.45), cw - Inches(0.6), Inches(0.5),
             lede, font=FONT_BODY, size=15, bold=True, color=color)
    add_text(s, x + Inches(0.3), y0 + Inches(2.0), cw - Inches(0.6), Inches(2.0),
             body, font=FONT_BODY, size=13, color=INK)


# ============================================================
# SLIDE 7 — The artifact (knowledge graph visual)
# ============================================================
s = prs.slides.add_slide(blank); add_bg(s)
add_chrome(s, "// 06 artifact", 7, TOTAL)
section_label(s, "The artifact")

add_text(s, Inches(0.6), Inches(1.35), Inches(11.5), Inches(0.8),
         "A growing, verified record.",
         font=FONT_BODY, size=40, bold=True, color=INK)

# left side text, right side simple graph visual
add_text(s, Inches(0.6), Inches(2.6), Inches(6.0), Inches(0.4),
         "// PROPERTIES", font=FONT_MONO, size=11, color=ACCENT, tracking=3)

props = [
    ("Free",       "to read for anyone in the world."),
    ("Verified",   "where machines can verify; reviewed where they can't."),
    ("Attributed", "to the researcher whose agent produced it."),
    ("Permanent",  "no walled garden, no platform lock-in."),
    ("Living",     "every accepted result expands the graph."),
]
yy = Inches(3.05)
for h, t in props:
    add_text(s, Inches(0.6), yy, Inches(2.0), Inches(0.35),
             h.upper(), font=FONT_MONO, size=12, bold=True, color=AMBER, tracking=2)
    add_text(s, Inches(2.5), yy, Inches(4.5), Inches(0.35),
             t, font=FONT_BODY, size=14, color=INK)
    yy += Inches(0.55)

# graph visual — nodes & edges on the right half
gx0, gy0 = Inches(7.4), Inches(2.5)
gw, gh   = Inches(5.3), Inches(4.0)
add_rect(s, gx0, gy0, gw, gh, PANEL, line=RULE)

# define some node positions (relative to gx0, gy0)
import random
random.seed(7)
nodes = [
    (1.0, 0.6,  ACCENT, "axiom"),
    (2.6, 1.2,  ACCENT, "lemma"),
    (4.2, 0.7,  ACCENT, "lemma"),
    (1.5, 2.4,  GREEN,  "thm"),
    (3.3, 2.6,  GREEN,  "thm"),
    (4.6, 2.0,  GREEN,  "thm"),
    (2.2, 3.4,  AMBER,  "conj"),
    (3.9, 3.5,  AMBER,  "conj"),
    (0.6, 1.8,  ACCENT, "lemma"),
    (4.7, 3.2,  GREEN,  "thm"),
]
edges = [(0,1),(1,2),(0,3),(1,4),(2,5),(3,4),(4,5),(3,6),(4,7),(5,7),(8,3),(8,0),(5,9),(7,9)]

def emu_pos(rel_x, rel_y):
    return gx0 + Inches(rel_x), gy0 + Inches(rel_y)

# draw edges first
for a, b in edges:
    ax, ay = emu_pos(nodes[a][0], nodes[a][1])
    bx, by = emu_pos(nodes[b][0], nodes[b][1])
    line = s.shapes.add_connector(1, ax, ay, bx, by)
    line.line.color.rgb = RULE
    line.line.width = Pt(0.75)

# draw nodes
for x, y, color, kind in nodes:
    nx, ny = emu_pos(x, y)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, nx - Emu(45720), ny - Emu(45720),
                             Emu(91440), Emu(91440))  # ~10pt circle
    dot.fill.solid(); dot.fill.fore_color.rgb = color
    dot.line.color.rgb = BG
    dot.line.width = Pt(1.0)
    dot.shadow.inherit = False

# legend
lgx, lgy = gx0 + Inches(0.25), gy0 + gh - Inches(0.5)
def legend_dot(x, color, label):
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, x, lgy + Inches(0.05),
                           Emu(76200), Emu(76200))
    d.fill.solid(); d.fill.fore_color.rgb = color
    d.line.fill.background(); d.shadow.inherit = False
    add_text(s, x + Inches(0.18), lgy, Inches(1.5), Inches(0.3),
             label, font=FONT_MONO, size=10, color=DIM, tracking=1)

legend_dot(lgx,                  ACCENT, "axiom / lemma")
legend_dot(lgx + Inches(1.7),    GREEN,  "verified theorem")
legend_dot(lgx + Inches(3.5),    AMBER,  "conjecture")


# ============================================================
# SLIDE 8 — The community math forgot it had
# ============================================================
s = prs.slides.add_slide(blank); add_bg(s)
add_chrome(s, "// 07 community", 8, TOTAL)
section_label(s, "The forum math lost")

add_text(s, Inches(0.6), Inches(1.35), Inches(11.5), Inches(0.8),
         "Rebuild the square AI emptied.",
         font=FONT_BODY, size=40, bold=True, color=INK)

# left half — the void
add_text(s, Inches(0.6), Inches(2.6), Inches(6.0), Inches(0.4),
         "// THE VOID", font=FONT_MONO, size=11, color=DIM, tracking=3)
add_text(s, Inches(0.6), Inches(3.0), Inches(6.0), Inches(3.4),
         "MathOverflow and Stack Exchange were the public square of working "
         "mathematics — where postdocs argued with experts, where a niche "
         "subfield's last living authority surfaced, where credit and reputation "
         "compounded.\n\n"
         "AI killed the question-and-answer half. The community half was "
         "collateral damage. Nothing has replaced it.",
         font=FONT_BODY, size=14, color=INK)

# right half — what we build
add_rect(s, Inches(7.1), Inches(2.5), Inches(5.6), Inches(4.0), PANEL, line=ACCENT)
add_text(s, Inches(7.4), Inches(2.7), Inches(5.0), Inches(0.4),
         "// WHAT WE BUILD", font=FONT_MONO, size=11, color=ACCENT, tracking=3)
add_text(s, Inches(7.4), Inches(3.1), Inches(5.0), Inches(0.55),
         "A community native\nto the agent era.",
         font=FONT_BODY, size=22, bold=True, color=INK)

bullets = [
    ("THREADS",     "Mathematicians discuss what agents produced — and didn't."),
    ("AGENTS",      "Share, fork, and remix the agents that did the work."),
    ("EXPERTS",     "The 80-year-old who holds a niche subfield gets a real role."),
    ("SIGNAL",      "Reputation compounds; attribution is permanent and public."),
]
yy = Inches(4.3)
for label, body in bullets:
    add_text(s, Inches(7.4), yy, Inches(1.5), Inches(0.3),
             label, font=FONT_MONO, size=10, bold=True, color=AMBER, tracking=2)
    add_text(s, Inches(8.7), yy, Inches(3.95), Inches(0.45),
             body, font=FONT_BODY, size=12, color=INK)
    yy += Inches(0.50)


# ============================================================
# SLIDE 9 — Verification is social
# ============================================================
s = prs.slides.add_slide(blank); add_bg(s)
add_chrome(s, "// 08 social proof", 9, TOTAL)
section_label(s, "Verification is social")

add_text(s, Inches(0.6), Inches(1.35), Inches(11.5), Inches(0.8),
         "Machines check. People bless.",
         font=FONT_BODY, size=40, bold=True, color=INK)

add_text(s, Inches(0.6), Inches(2.55), Inches(11.5), Inches(0.5),
         "Trust is layered. Every claim agents produce moves through stacked "
         "checkmarks — none of them faked, all of them visible.",
         font=FONT_BODY, size=15, color=DIM)

# four check-mark badges in a row
badges = [
    ("MACHINE",    "Lean-verified",   "Mechanically certain.",   ACCENT),
    ("EXPERT",     "Domain-reviewed", "Blessed by an authority on the subfield.", GREEN),
    ("COMMUNITY",  "Peer-discussed",  "Survived open thread review.", AMBER),
    ("REPUTATION", "Author-attributed","The mathematician who shaped the agent owns it.", MAGENTA),
]
y0 = Inches(3.55); cw = Inches(2.95); ch = Inches(2.85); gap = Inches(0.16)
for i, (kind, head, body, color) in enumerate(badges):
    x = Inches(0.6) + i * (cw + gap)
    add_rect(s, x, y0, cw, ch, PANEL, line=RULE)
    # checkmark glyph (drawn with two short connectors)
    cx = x + Inches(0.4); cy = y0 + Inches(0.5)
    add_rect(s, cx, cy, Inches(0.45), Inches(0.45), color)
    add_text(s, cx, cy - Inches(0.04), Inches(0.45), Inches(0.5),
             "✓", font=FONT_BODY, size=22, bold=True, color=BG,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(1.0), y0 + Inches(0.45), cw - Inches(1.2), Inches(0.4),
             kind, font=FONT_MONO, size=10, color=color, tracking=2)
    add_text(s, x + Inches(1.0), y0 + Inches(0.78), cw - Inches(1.2), Inches(0.5),
             head, font=FONT_BODY, size=16, bold=True, color=INK)
    add_text(s, x + Inches(0.3), y0 + Inches(1.65), cw - Inches(0.6), Inches(1.1),
             body, font=FONT_BODY, size=13, color=DIM)


# ============================================================
# SLIDE 10 — Why mathematicians will adopt it
# ============================================================
s = prs.slides.add_slide(blank); add_bg(s)
add_chrome(s, "// 09 adoption", 10, TOTAL)
section_label(s, "Why they will adopt")

add_text(s, Inches(0.6), Inches(1.35), Inches(11.5), Inches(0.8),
         "Built for how mathematicians actually work.",
         font=FONT_BODY, size=36, bold=True, color=INK)

rows = [
    ("TRUST",       "Mechanical verification. No flashy claims. No marketing-grade math."),
    ("ATTRIBUTION", "Their name on their results — by default, permanently, citably."),
    ("PERSISTENCE", "An agent that remembers, accumulates, and resumes. No more sandcastles."),
    ("AGENCY",      "Their agent, their tools, their prompts. Workflow, not chatbot."),
    ("AESTHETICS",  "A platform that looks and feels like it was built by people who do math."),
]
yy = Inches(2.7)
for label, body in rows:
    add_rect(s, Inches(0.6), yy, Inches(12.13), Emu(9525), RULE)
    add_text(s, Inches(0.6), yy + Inches(0.18), Inches(2.6), Inches(0.4),
             label, font=FONT_MONO, size=14, bold=True, color=ACCENT, tracking=3)
    add_text(s, Inches(3.4), yy + Inches(0.18), Inches(9.3), Inches(0.4),
             body, font=FONT_BODY, size=18, color=INK)
    yy += Inches(0.78)


# ============================================================
# SLIDE 9 — Beyond math
# ============================================================
s = prs.slides.add_slide(blank); add_bg(s)
add_chrome(s, "// 10 horizon", 11, TOTAL)
section_label(s, "Horizon")

add_text(s, Inches(0.6), Inches(1.35), Inches(11.5), Inches(0.8),
         "Math is the wedge.",
         font=FONT_BODY, size=48, bold=True, color=INK)
add_text(s, Inches(0.6), Inches(2.4), Inches(11.5), Inches(0.6),
         "The architecture is universal.",
         font=FONT_BODY, size=28, color=DIM)

# horizontal arrow of fields
fields = [
    ("MATHEMATICS",    "the wedge",    ACCENT,  True),
    ("PHYSICS",        "next neighbor", AMBER,  False),
    ("CHEMISTRY",      "right pattern", AMBER,  False),
    ("ANY DOMAIN",     "where reasoning produces results", DIM, False),
]
y0 = Inches(4.0); cw = Inches(2.95); ch = Inches(2.4); gap = Inches(0.16)
for i, (h, sub, color, filled) in enumerate(fields):
    x = Inches(0.6) + i * (cw + gap)
    add_rect(s, x, y0, cw, ch, PANEL if filled else BG,
             line=RULE if not filled else None)
    if filled:
        add_rect(s, x, y0, cw, Inches(0.08), color)
    add_text(s, x + Inches(0.3), y0 + Inches(0.4), cw - Inches(0.6), Inches(0.4),
             f"0{i+1}", font=FONT_MONO, size=10, color=color, tracking=2)
    add_text(s, x + Inches(0.3), y0 + Inches(0.75), cw - Inches(0.6), Inches(0.6),
             h, font=FONT_BODY, size=20, bold=True, color=INK, tracking=1)
    add_text(s, x + Inches(0.3), y0 + Inches(1.5), cw - Inches(0.6), Inches(0.7),
             sub, font=FONT_MONO, size=11, color=color, tracking=1)


# ============================================================
# SLIDE 10 — Culture / aesthetic
# ============================================================
s = prs.slides.add_slide(blank); add_bg(s)
add_chrome(s, "// 11 culture", 12, TOTAL)
section_label(s, "Culture")

add_text(s, Inches(0.6), Inches(1.35), Inches(11.5), Inches(0.8),
         "By mathematicians, for mathematicians.",
         font=FONT_BODY, size=36, bold=True, color=INK)

quote_y = Inches(3.0)
# big accent bar
add_rect(s, Inches(0.6), quote_y, Emu(76200), Inches(2.5), AMBER)
add_text(s, Inches(1.0), quote_y - Inches(0.05), Inches(11), Inches(2.6),
         "“Mathematicians pay dearly\nfor instruments they trust\nand ignore everything else.”",
         font=FONT_BODY, size=34, bold=True, color=INK)
add_text(s, Inches(1.0), quote_y + Inches(2.45), Inches(11), Inches(0.4),
         "// founding principle",
         font=FONT_MONO, size=11, color=AMBER, tracking=3)


# ============================================================
# SLIDE 11 — Closing
# ============================================================
s = prs.slides.add_slide(blank); add_bg(s)
add_chrome(s, "// 12 close", 13, TOTAL)

# big mark
add_rect(s, Inches(0.6), Inches(2.7), Emu(76200), Inches(2.4), ACCENT)
add_text(s, Inches(0.95), Inches(2.55), Inches(8), Inches(0.4),
         "MATHAI", font=FONT_MONO, size=12, color=ACCENT, tracking=6)
add_text(s, Inches(0.95), Inches(2.95), Inches(11.5), Inches(1.6),
         "We are building it now.",
         font=FONT_BODY, size=64, bold=True, color=INK)
add_text(s, Inches(0.95), Inches(4.6), Inches(11.5), Inches(0.6),
         "The map of mathematics, machine-verified, attributed to its makers,\n"
         "free to the world. Built by the people who actually do the work.",
         font=FONT_BODY, size=20, color=DIM)
add_text(s, Inches(0.95), Inches(6.0), Inches(11.5), Inches(0.4),
         "Anton Shakov  ·  Christian Kudeba",
         font=FONT_MONO, size=11, color=AMBER, tracking=2)


out = r"C:\Users\anton\OneDrive\Desktop\mathAI\mathAI_vision_deck.pptx"
prs.save(out)
print(f"wrote {out}")
print(f"slides: {len(prs.slides)}")
